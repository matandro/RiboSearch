from pptx import Presentation
from pptx.util import Inches
import os


BASE_FOLDER = 'D:\\Matan\\Dropbox\\PHd\\RiboSearch\\25_2'
MFE_FOLDER = os.path.join(BASE_FOLDER, 'MFE_seq_images')
CENTROID_FOLDER = os.path.join(BASE_FOLDER, 'centroid_seq_images')
SEQUENCE_FILE = os.path.join(BASE_FOLDER, 'possible_seq.txt')


def add_map_images(sequence_map):
    res_full_map = {}
    for key, value in sequence_map.items():
        image_name = "seq_{}.jpg".format(key)
        res_full_map[key] = {'sequence': value,
                             'mfe': os.path.join(MFE_FOLDER, image_name),
                             'centroid': os.path.join(CENTROID_FOLDER, image_name)}
    return res_full_map


def generate_seq_map():
    res_seq_map = {}
    with open(SEQUENCE_FILE, 'r') as seq_file:
        seq_index = 1
        for line in seq_file:
            res_seq_map[seq_index] = line.strip()
            seq_index += 1
    return res_seq_map


seq_map = generate_seq_map()
full_map = add_map_images(seq_map)

prs = Presentation()
# default slide width
prs.slide_width = 9144000
# slide height @ 16:9
prs.slide_height = 5143500
blank_slide_layout = prs.slide_layouts[6]

for i in range(1, len(full_map) + 1):
    single_info = full_map.get(i)
    slide = prs.slides.add_slide(blank_slide_layout)
    # set pic top
    pic_top = prs.slide_height * 0.10
    item_width = prs.slide_width * 0.45
    # add sequence as title
    title_shape = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
    title_shape.text = single_info.get('sequence')
    # add MFE image
    left = prs.slide_width * 0.04
    slide.shapes.add_picture(single_info.get('mfe'), left, pic_top,
                             width=item_width, height=prs.slide_height * 0.85)
    mfe_text = slide.shapes.add_textbox(left, prs.slide_height * 0.95,
                                        item_width, prs.slide_height * 0.03)
    mfe_text.text = 'MFE structure'
    # add centroid image
    left = prs.slide_width * 0.51
    slide.shapes.add_picture(single_info.get('centroid'), left, pic_top,
                             width=item_width, height=prs.slide_height * 0.85)
    centroid_text = slide.shapes.add_textbox(left, prs.slide_height * 0.95,
                                             item_width, prs.slide_height * 0.03)
    centroid_text.text = 'Centroid structure'
    # empty slide after each
    slide = prs.slides.add_slide(blank_slide_layout)
prs.save(os.path.join(BASE_FOLDER, 'report_frame.pptx'))
